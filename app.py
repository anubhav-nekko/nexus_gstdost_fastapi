import os
import io
import json
import pickle
import tempfile
from typing import List, Dict, Any, Tuple, Optional

import boto3 # type: ignore
import faiss # type: ignore
import fitz # type: ignore
import numpy as np # type: ignore
import pandas as pd # type: ignore
import requests # type: ignore
from fastapi import FastAPI, File, UploadFile, HTTPException, Body # type: ignore
from fastapi.responses import JSONResponse # type: ignore
from pydantic import BaseModel # type: ignore
from sentence_transformers import SentenceTransformer # type: ignore
from tavily import TavilyClient # type: ignore
from docx import Document # type: ignore
from pptx import Presentation # type: ignore
from PIL import Image # type: ignore
from openai import AzureOpenAI, OpenAI # type: ignore

###############################################################################
# 1) CONFIGURATION & PROMPT LIBRARIES
###############################################################################

# =============== Load secrets from JSON (adjust path as needed) ==============
SECRETS_PATH = "../secrets.json"

def load_secrets(path: str) -> dict:
    with open(path, "r", encoding="utf-8") as f:
        return json.load(f)

SECRETS = load_secrets(SECRETS_PATH)

aws_access_key_id = SECRETS["aws_access_key_id"]
aws_secret_access_key = SECRETS["aws_secret_access_key"]

# Note: You have separate ARNs for Titan / Claude / Deepseek:
INFERENCE_PROFILE_ARN = SECRETS["INFERENCE_PROFILE_ARN"]
INFERENCE_PROFILE_CLAUDE = SECRETS["INFERENCE_PROFILE_CLAUDE"]
INFERENCE_PROFILE_DEEPSEEK = SECRETS["INFERENCE_PROFILE_DEEPSEEK"]
REGION = SECRETS["REGION"]
REGION2 = SECRETS["REGION2"]

GPT_ENDPOINT = SECRETS["GPT_ENDPOINT"]
GPT_API = SECRETS["GPT_API"]
OPENAI_KEY = SECRETS["OPENAI_KEY"]

TAVILY_API = SECRETS["TAVILY_API"]

# S3 Bucket Name & Local paths to FAISS index
S3_BUCKET_NAME = SECRETS["s3_bucket_name"]
FAISS_INDEX_PATH = SECRETS["FAISS_INDEX_PATH"]
METADATA_STORE_PATH = SECRETS["METADATA_STORE_PATH"]

# Clients for AWS
s3_client = boto3.client(
    "s3",
    region_name=REGION,
    aws_access_key_id=aws_access_key_id,
    aws_secret_access_key=aws_secret_access_key
)
textract_client = boto3.client(
    "textract",
    region_name=REGION,
    aws_access_key_id=aws_access_key_id,
    aws_secret_access_key=aws_secret_access_key
)

# Possibly separate client for bedrock in REGION2
bedrock_runtime = boto3.client(
    "bedrock-runtime",
    region_name=REGION2,
    aws_access_key_id=aws_access_key_id,
    aws_secret_access_key=aws_secret_access_key
)

# =============== SentenceTransformer MPNet for embeddings =====================
mpnet_model = SentenceTransformer("sentence-transformers/all-mpnet-base-v2")

# =============== Initialize FAISS index & metadata in memory ==================
DIMENSION = 768  # for mpnet base v2
faiss_index = faiss.IndexFlatL2(DIMENSION)
metadata_store = []

# =============== Some helpful global prompts (if needed) ======================
prompt_library = {
    "custom": "",
    "Summarization": "You are an advanced legal AI designed to summarize allegations ... [Truncated for brevity]",
    "Chronological Event Extraction": "...",
    "Disputed Amount Details": "...",
    "Relevant Legal Framework Identification": "...",
    "Taxpayer Argument": "...",
}

system_message = """
You are an advanced legal data analyst specializing in legal document analysis. Provide an in-depth analysis...
"""

summary_prompt = "You are a Helpful Legal Data Analyst specializing in tax-related legal document analysis..."
insights_prompt = "You are a Helpful Legal Data Analyst specializing in tax-related legal document analysis..."
ws_prompt_lib = "You are a legal research assistant tasked with compiling relevant legal cases..."
insights_prompt_lib = "You are a legal research and analysis assistant..."
qna_prompt = "You are a legal research and analysis assistant..."
nekkollm_prompt = "You are Nekko LLM..."

###############################################################################
# 2) FASTAPI APP
###############################################################################
app = FastAPI(title="Document Query Assistant", version="1.0.0")

###############################################################################
# 3) UTILITY FUNCTIONS: S3, FAISS, Textract, etc.
###############################################################################

def file_exists_in_s3(key: str) -> bool:
    """Check if an object with key `key` exists in S3."""
    try:
        s3_client.head_object(Bucket=S3_BUCKET_NAME, Key=key)
        return True
    except Exception as e:
        if hasattr(e, "response") and e.response["Error"]["Code"] == "404":
            return False
        raise e

def upload_to_s3(local_path: str, s3_key: str):
    """Upload local file to S3 bucket under `s3_key`."""
    try:
        s3_client.upload_file(local_path, S3_BUCKET_NAME, s3_key)
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Failed to upload {local_path} -> S3: {e}")

def download_from_s3(s3_key: str, local_path: str) -> bool:
    """Download file from S3. Return True if successful, False if not found."""
    try:
        s3_client.download_file(S3_BUCKET_NAME, s3_key, local_path)
        return True
    except Exception as e:
        if hasattr(e, "response") and e.response["Error"]["Code"] == "404":
            return False
        raise HTTPException(status_code=500, detail=f"Failed to download {s3_key} from S3: {e}")

def generate_embeddings(text: str) -> np.ndarray:
    """Generate MPNet embeddings for the given text."""
    emb = mpnet_model.encode(text, normalize_embeddings=True)
    return emb.astype(np.float32)

def read_index_and_metadata():
    """Load the FAISS index and metadata store from S3 if available."""
    global faiss_index, metadata_store

    # Download both the index and metadata
    index_filename = os.path.basename(FAISS_INDEX_PATH)
    meta_filename = os.path.basename(METADATA_STORE_PATH)

    index_ok = download_from_s3(index_filename, FAISS_INDEX_PATH)
    meta_ok = download_from_s3(meta_filename, METADATA_STORE_PATH)

    if index_ok and os.path.exists(FAISS_INDEX_PATH):
        try:
            faiss_index = faiss.read_index(FAISS_INDEX_PATH)
        except Exception as e:
            print("Error reading index from disk:", e)
            faiss_index = faiss.IndexFlatL2(DIMENSION)
    else:
        faiss_index = faiss.IndexFlatL2(DIMENSION)

    if meta_ok and os.path.exists(METADATA_STORE_PATH):
        try:
            with open(METADATA_STORE_PATH, "rb") as f:
                metadata_store.clear()
                metadata_store.extend(pickle.load(f))
        except Exception as e:
            print("Error reading metadata from disk:", e)
            metadata_store.clear()
    else:
        metadata_store.clear()

def save_index_and_metadata():
    """Save FAISS index + metadata to local disk, then upload to S3."""
    # Save index
    faiss.write_index(faiss_index, FAISS_INDEX_PATH)
    # Save metadata
    with open(METADATA_STORE_PATH, "wb") as f:
        pickle.dump(metadata_store, f)

    # Upload to S3
    index_filename = os.path.basename(FAISS_INDEX_PATH)
    meta_filename = os.path.basename(METADATA_STORE_PATH)
    upload_to_s3(FAISS_INDEX_PATH, index_filename)
    upload_to_s3(METADATA_STORE_PATH, meta_filename)


@app.on_event("startup")
def on_start():
    """Load index/metadata from S3 on startup."""
    read_index_and_metadata()
    print("Index and metadata loaded.")

def textract_pdf_to_text_pages(pdf_path: str) -> List[Tuple[int, str]]:
    """
    For each page in the PDF, produce (page_num, extracted_text).
    Uses AWS Textract detect_document_text on page images.
    """
    doc = fitz.open(pdf_path)
    results = []
    for i in range(len(doc)):
        page = doc.load_page(i)
        pix = page.get_pixmap()
        temp_png = os.path.join(tempfile.gettempdir(), f"page_{i}.png")
        pix.save(temp_png)

        with open(temp_png, "rb") as f:
            data = f.read()

        try:
            response = textract_client.detect_document_text(Document={"Bytes": data})
            lines = []
            for block in response.get("Blocks", []):
                if block["BlockType"] == "LINE" and "Text" in block:
                    lines.append(block["Text"])
            page_text = "\n".join(lines)
        except Exception as e:
            page_text = f"[ERROR reading page {i+1}: {str(e)}]"
        finally:
            if os.path.exists(temp_png):
                os.remove(temp_png)

        results.append((i + 1, page_text))
    return results

def docx_to_chunks(docx_path: str) -> List[Tuple[int, str]]:
    """Read docx content in ~1000-char chunks. Return (chunk_index, chunk_text)."""
    document = Document(docx_path)
    full_text = "\n".join([p.text for p in document.paragraphs if p.text.strip()])
    chunks = []
    chunk_size = 1000
    i = 0
    start = 0
    while start < len(full_text):
        end = min(start + chunk_size, len(full_text))
        chunk = full_text[start:end]
        i += 1
        chunks.append((i, chunk))
        start = end
    return chunks

def pptx_to_chunks(pptx_path: str) -> List[Tuple[int, str]]:
    """Extract text from slides. Return (slide_num, text)."""
    prs = Presentation(pptx_path)
    results = []
    for idx, slide in enumerate(prs.slides, 1):
        lines = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                lines.append(shape.text)
        results.append((idx, "\n".join(lines)))
    return results

def xlsx_to_chunks(xlsx_path: str) -> List[Tuple[int, str]]:
    """
    Reads first sheet. Splits rows in chunks of 50 for embedding. 
    Return (chunk_index, text_of_that_chunk).
    """
    # If CSV, handle with pd.read_csv
    ext = os.path.splitext(xlsx_path)[1].lower()
    if ext == ".csv":
        df = pd.read_csv(xlsx_path)
    else:
        xls = pd.ExcelFile(xlsx_path)
        first_sheet = xls.sheet_names[0]
        df = pd.read_excel(xls, sheet_name=first_sheet)

    chunk_size = 50
    rows = len(df)
    results = []
    idx = 0
    for start in range(0, rows, chunk_size):
        end = start + chunk_size
        chunk_df = df.iloc[start:end]
        text_block = chunk_df.to_string(index=False)
        idx += 1
        results.append((idx, text_block))
    return results

def image_to_text(image_path: str) -> str:
    """OCR using Textract for a single image."""
    with open(image_path, "rb") as f:
        data = f.read()
    response = textract_client.detect_document_text(Document={"Bytes": data})
    lines = []
    for block in response.get("Blocks", []):
        if block["BlockType"] == "LINE" and "Text" in block:
            lines.append(block["Text"])
    return "\n".join(lines)

###############################################################################
# 4) LLM INVOCATIONS
###############################################################################

def call_llm_bedrock_titan(system_msg: str, user_msg: str) -> str:
    # Combine system and user messages
    messages = system_msg + user_msg

    # Prepare the request payload
    payload = {
        "anthropic_version": "bedrock-2023-05-31",
        "max_tokens": 4096,
        "messages": [
            {
                "role": "user",
                "content": messages
            }
        ]
    }

    try:
        # Invoke the model (Claude)
        response = bedrock_runtime.invoke_model(
            modelId=INFERENCE_PROFILE_ARN,  # Use the ARN for your inference profile
            contentType='application/json',
            accept='application/json',
            body=json.dumps(payload)
        )

        # Parse and return the response
        response_body = json.loads(response['body'].read())
        return response_body['content'][0]['text']

    except Exception as e:
        return f"An error occurred: {str(e)}"

def call_claude_35(system_msg: str, user_msg: str) -> str:
    # Combine system and user messages
    messages = system_msg + user_msg

    # Prepare the request payload
    payload = {
        "anthropic_version": "bedrock-2023-05-31",
        "max_tokens": 4096,
        "messages": [
            {
                "role": "user",
                "content": messages
            }
        ]
    }

    try:
        # Invoke the model (Claude)
        response = bedrock_runtime.invoke_model(
            modelId=INFERENCE_PROFILE_ARN,  # Use the ARN for your inference profile
            contentType='application/json',
            accept='application/json',
            body=json.dumps(payload)
        )

        # Parse and return the response
        response_body = json.loads(response['body'].read())
        return response_body['content'][0]['text']

    except Exception as e:
        return f"An error occurred: {str(e)}"

def call_claude_37(system_msg: str, user_msg: str) -> str:
    # Combine system and user messages
    messages = system_msg + user_msg

    # Prepare the request payload
    payload = {
        "anthropic_version": "bedrock-2023-05-31",
        "max_tokens": 4096,
        "messages": [
            {
                "role": "user",
                "content": messages
            }
        ]
    }

    try:
        # Invoke the model (Claude)
        response = bedrock_runtime.invoke_model(
            modelId=INFERENCE_PROFILE_CLAUDE,  # Use the ARN for your inference profile
            contentType='application/json',
            accept='application/json',
            body=json.dumps(payload)
        )

        # Parse and return the response
        response_body = json.loads(response['body'].read())
        return response_body['content'][0]['text']

    except Exception as e:
        return f"An error occurred: {str(e)}"

def call_deepseek(system_msg: str, user_msg: str) -> str:
    # Combine system and user messages
    messages = system_msg + user_msg

    # Prepare the request payload
    payload = {
        "max_tokens": 4096,
        "messages": [
            {
                "role": "user",
                "content": messages
            }
        ]
    }

    try:
        # Invoke the model (Claude)
        response = bedrock_runtime.invoke_model(
            modelId=INFERENCE_PROFILE_DEEPSEEK,  # Use the ARN for your inference profile
            contentType='application/json',
            accept='application/json',
            body=json.dumps(payload)
        )

        # Parse and return the response
        response_body = json.loads(response['body'].read())
        return response_body["choices"][0]["message"]["content"]

    except Exception as e:
        return f"An error occurred: {str(e)}"

def call_gpt_4o(system_msg: str, user_msg: str) -> str:
    # url = GPT_ENDPOINT
    # headers = {  
    #     "Content-Type": "application/json",  
    #     "api-key": GPT_API
    # }  
    # messages = [
    #     {"role": "system", "content": system_message},
    #     {"role": "user", "content": user_query}
    # ]
    # payload = {  
    #     "messages": messages,  
    #     "temperature": 0.7,  
    #     "max_tokens": 4096   
    # }
    # response = requests.post(url, headers=headers, data=json.dumps(payload))
    # response.raise_for_status()  
    # return response.json()["choices"][0]["message"]["content"]
    
    client = OpenAI(api_key = OPENAI_KEY)
    
    response = client.responses.create(
      model="gpt-4o",
      input=[
        {
          "role": "system",
          "content": [
            {
              "type": "input_text",
              "text": system_msg
            }
          ]
        },
        {
          "role": "user",
          "content": [
            {
              "type": "input_text",
              "text": user_msg
            }
          ]
        }
      ],
      text={
        "format": {
          "type": "text"
        }
      },
      reasoning={},
      tools=[],
      temperature=0.7,
      max_output_tokens=4096,
      top_p=1,
      store=True
    )
    
    return response.output[0].content[0].text

def call_llm_by_name(model_name: str, system_msg: str, user_msg: str) -> str:
    """
    Dispatch to the correct LLM function based on `model_name`.
    """
    model_name_lower = model_name.strip().lower()
    if model_name_lower == "claude 3.5 sonnet":
        return call_claude_35(system_msg, user_msg)
    elif model_name_lower == "claude 3.7 sonnet":
        return call_claude_37(system_msg, user_msg)
    elif model_name_lower == "deepseek r1":
        return call_deepseek(system_msg, user_msg)
    elif model_name_lower == "gpt 4o":
        return call_gpt_4o(system_msg, user_msg)
    else:
        # fallback to default
        return call_llm_bedrock_titan(system_msg, user_msg)

###############################################################################
# 5) REQUEST MODELS
###############################################################################

class UploadResponse(BaseModel):
    message: str
    filename: str

class QueryRequest(BaseModel):
    selected_files: List[str]
    selected_page_ranges: Dict[str, Tuple[int,int]]  # {"filename.pdf": [1, 5], ...}
    prompt: str
    top_k: int
    last_messages: List[str] = []
    web_search: bool = False
    llm_model: str = "Claude 3.5 Sonnet"
    draft_mode: bool = False
    analyse_mode: bool = False

class QueryResponse(BaseModel):
    answer: str
    top_k_metadata: List[Any]
    web_search_results: Any = None

###############################################################################
# 6) UPLOAD DOCUMENTS (ALL FILE TYPES) ENDPOINT
###############################################################################

@app.post("/upload_document", response_model=UploadResponse)
def upload_document(file: UploadFile = File(...)):
    """
    Upload a single document (PDF, image, DOCX, PPTX, XLSX, CSV).
    - Saves to a temp file.
    - Uploads to S3 if not already present.
    - Extracts text, embeddings it in FAISS, stores in metadata.
    """
    try:
        filename = file.filename
        ext = os.path.splitext(filename)[1].lower()

        # 1) Save temp
        temp_path = os.path.join(tempfile.gettempdir(), filename)
        with open(temp_path, "wb") as f:
            f.write(file.file.read())

        # 2) Check if file in S3
        if file_exists_in_s3(filename):
            return UploadResponse(
                message=f"File {filename} already in S3; skipping re-upload.",
                filename=filename
            )
        else:
            # Upload to S3
            upload_to_s3(temp_path, filename)

        # 3) Extract text from the file, chunk if needed, embed
        new_records = []
        if ext == ".pdf":
            pages = textract_pdf_to_text_pages(temp_path)
            for (page_num, txt) in pages:
                if txt.strip():
                    emb = generate_embeddings(txt)
                    faiss_index.add(emb.reshape(1, -1))
                    new_records.append({
                        "filename": filename,
                        "page": page_num,
                        "text": txt
                    })
        elif ext in [".jpg", ".jpeg", ".png"]:
            # treat as single-page image
            txt = image_to_text(temp_path)
            emb = generate_embeddings(txt)
            faiss_index.add(emb.reshape(1, -1))
            new_records.append({
                "filename": filename,
                "page": 1,
                "text": txt
            })
        elif ext in [".doc", ".docx"]:
            chunks = docx_to_chunks(temp_path)
            for (idx, chunk_text) in chunks:
                if chunk_text.strip():
                    emb = generate_embeddings(chunk_text)
                    faiss_index.add(emb.reshape(1, -1))
                    new_records.append({
                        "filename": filename,
                        "page": idx,
                        "text": chunk_text
                    })
        elif ext == ".pptx":
            slides = pptx_to_chunks(temp_path)
            for (idx, slide_text) in slides:
                if slide_text.strip():
                    emb = generate_embeddings(slide_text)
                    faiss_index.add(emb.reshape(1, -1))
                    new_records.append({
                        "filename": filename,
                        "page": idx,
                        "text": slide_text
                    })
        elif ext in [".xlsx", ".csv"]:
            table_chunks = xlsx_to_chunks(temp_path)
            for (idx, chunk_text) in table_chunks:
                if chunk_text.strip():
                    emb = generate_embeddings(chunk_text)
                    faiss_index.add(emb.reshape(1, -1))
                    new_records.append({
                        "filename": filename,
                        "page": idx,
                        "text": chunk_text
                    })
        else:
            raise HTTPException(status_code=400, detail=f"Unsupported file type {ext}")

        # 4) Add to metadata store
        metadata_store.extend(new_records)

        # 5) Save the index & metadata
        save_index_and_metadata()

        # 6) Cleanup
        if os.path.exists(temp_path):
            os.remove(temp_path)

        return UploadResponse(
            message=f"File {filename} processed and indexed successfully.",
            filename=filename
        )

    except Exception as ex:
        raise HTTPException(status_code=500, detail=str(ex))

###############################################################################
# 7) QUERY DOCUMENTS WITH PAGE RANGE
###############################################################################

@app.post("/query_documents_with_page_range", response_model=QueryResponse)
def query_documents(req: QueryRequest):
    """
    1) Takes a user prompt plus user-specified selected_files, selected_page_ranges, top_k, etc.
    2) Optionally refines the prompt using an "intelligent query refiner" (like your code).
    3) Conducts semantic search in FAISS.
    4) Filters results by file + page range.
    5) If web_search is True, calls Távily with the 'web_search_prompt'.
    6) If draft_mode or analyse_mode is True, incorporate that into final user prompt logic.
    7) Calls the chosen LLM (Claude 3.5, 3.7, Deepseek R1, or GPT 4o) with the final prompt.
    8) Returns the final answer, the top_k metadata, and web search results.
    """
    # Step 0: If the index is empty, return
    if faiss_index.ntotal == 0:
        return QueryResponse(
            answer="No data available to query. The FAISS index is empty.",
            top_k_metadata=[],
            web_search_results=None
        )

    # Step 1: (Optional) refine the prompt for semantic search & web search
    # Here is a minimal approach that calls the default LLM to parse out "semantic_search_prompt" & "web_search_prompt".
    refine_system_msg = (
        "You are an intelligent query refiner. "
        "Take the user's original query + last few messages and produce JSON with keys: 'semantic_search_prompt' and 'web_search_prompt'."
    )
    last_msgs_concat = "\n".join(req.last_messages)
    refine_user_msg = f"User Query: {req.prompt}\n\nLast Messages: {last_msgs_concat}\n\nGenerate the JSON."
    refinement_raw = call_llm_bedrock_titan(refine_system_msg, refine_user_msg)

    # Attempt to parse
    semantic_search_prompt = req.prompt
    web_search_prompt = req.prompt
    try:
        # Typically you'd do something like:
        # raw_json = refinement_raw.split("```json")[1].split("```")[0]
        # But for safety:
        parts = refinement_raw.split("```json")
        if len(parts) > 1:
            jstring = parts[1].split("```")[0]
            parsed = json.loads(jstring)
            semantic_search_prompt = parsed.get("semantic_search_prompt", req.prompt)
            web_search_prompt = parsed.get("web_search_prompt", req.prompt)
    except:
        pass

    # Step 2: Generate embedding for the semantic_search_prompt
    query_emb = generate_embeddings(semantic_search_prompt).reshape(1, -1)
    k = faiss_index.ntotal
    distances, indices = faiss_index.search(query_emb, k)

    # Step 3: Filter results for chosen files + page ranges
    filtered = []
    for dist, idx in zip(distances[0], indices[0]):
        if idx < len(metadata_store):
            record = metadata_store[idx]
            if record["filename"] in req.selected_files:
                (pg_start, pg_end) = req.selected_page_ranges.get(record["filename"], (1, 999999))
                if pg_start <= record["page"] <= pg_end:
                    filtered.append((dist, idx))

    # sort by distance ascending, take top_k
    topk = sorted(filtered, key=lambda x: x[0])[: req.top_k]
    top_k_metadata = [metadata_store[i] for (_, i) in topk]

    # Step 4: Web search if requested
    web_search_results = None
    if req.web_search:
        try:
            tavily_client = TavilyClient(api_key=TAVILY_API)
            web_search_results = tavily_client.search(
                query=web_search_prompt,
                search_depth="advanced",
                include_raw_content=True
            )
        except Exception as e:
            web_search_results = {"error": str(e)}

    # Step 5: Build final user prompt
    # If draft_mode or analyse_mode, we do additional steps. We'll do a simple approach:
    # We can just append a small note to the user prompt if these toggles are set:
    mode_context = ""
    if req.draft_mode:
        mode_context += "\n\n[Draft Mode Active: The user wants a drafted legal text/argument form.]"
    if req.analyse_mode:
        mode_context += "\n\n[Analyse Mode Active: Provide deeper, more thorough analysis.]"

    # Build final
    context_blob = json.dumps(top_k_metadata, indent=2)
    final_user_prompt = (
        f"User's Query:\n{req.prompt}\n"
        f"\nRetrieved Document Context:\n{context_blob}\n"
        f"{mode_context}\n"
        "Please provide a detailed, step-by-step answer."
    )

    # Step 6: Call the chosen LLM
    answer = call_llm_by_name(req.llm_model, system_message, final_user_prompt)

    return QueryResponse(
        answer=answer,
        top_k_metadata=top_k_metadata,
        web_search_results=web_search_results
    )
